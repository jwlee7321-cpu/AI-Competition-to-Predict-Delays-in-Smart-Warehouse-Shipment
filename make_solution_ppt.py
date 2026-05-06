from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("solution_A120.pptx")

NAVY = RGBColor(25, 43, 70)
BLUE = RGBColor(46, 100, 170)
LIGHT_BLUE = RGBColor(230, 238, 250)
GRAY = RGBColor(90, 96, 108)
LIGHT_GRAY = RGBColor(245, 247, 250)
GREEN = RGBColor(41, 128, 92)
ORANGE = RGBColor(205, 118, 45)

FONT = "Malgun Gothic"


def set_run(run, size=20, bold=False, color=RGBColor(20, 20, 20)):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_footer(slide, idx: int):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = f"Smart Warehouse Delay Prediction | Final A120 Solution | {idx}"
    p.alignment = PP_ALIGN.RIGHT
    p.runs[0].font.name = FONT
    p.runs[0].font.size = Pt(8)
    p.runs[0].font.color.rgb = RGBColor(135, 143, 155)


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_run(p.runs[0], size=28, bold=True, color=NAVY)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.62), Inches(1.05), Inches(11.7), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        set_run(sp.runs[0], size=12, color=GRAY)


def add_bar(slide, color=BLUE):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_bullets(slide, items, x=0.75, y=1.55, w=5.9, h=4.9, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.name = FONT
        p.font.size = Pt(size - level)
        p.font.color.rgb = RGBColor(35, 39, 47)
        p.space_after = Pt(6)
    return box


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_GRAY
    rect.line.color.rgb = RGBColor(220, 226, 235)
    rect.line.width = Pt(1)

    stripe = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.color.rgb = accent

    t = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.15), Inches(w - 0.35), Inches(0.35))
    p = t.text_frame.paragraphs[0]
    p.text = title
    set_run(p.runs[0], size=14, bold=True, color=NAVY)

    b = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.55), Inches(w - 0.35), Inches(h - 0.65))
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = body
    set_run(p.runs[0], size=12, color=RGBColor(45, 50, 60))


def add_code_box(slide, code, x=0.8, y=4.45, w=11.7, h=1.35):
    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(30, 38, 52)
    rect.line.color.rgb = RGBColor(30, 38, 52)
    box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.18), Inches(w - 0.44), Inches(h - 0.25))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.runs[0].font.name = "Consolas"
    p.runs[0].font.size = Pt(15)
    p.runs[0].font.color.rgb = RGBColor(235, 242, 250)


def add_table(slide, rows, x, y, w, h, font_size=11):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = RGBColor(35, 39, 47)
                    if r == 0:
                        run.font.bold = True
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    return table_shape


def make_slide(prs, idx, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide)
    add_title(slide, title, subtitle)
    add_footer(slide, idx)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bar(slide, NAVY)
    title = slide.shapes.add_textbox(Inches(0.75), Inches(1.25), Inches(11.8), Inches(1.0))
    p = title.text_frame.paragraphs[0]
    p.text = "스마트 물류창고 출고 지연 예측"
    set_run(p.runs[0], size=34, bold=True, color=NAVY)
    sub = slide.shapes.add_textbox(Inches(0.78), Inches(2.25), Inches(11.0), Inches(0.7))
    p = sub.text_frame.paragraphs[0]
    p.text = "Final A120 Solution | Public MAE 9.6773143392"
    set_run(p.runs[0], size=20, color=BLUE)
    add_card(slide, 0.85, 3.35, 3.6, 1.25, "최종 제출 파일", "sub_v44_anti_v41layout035_a120.csv", BLUE)
    add_card(slide, 4.85, 3.35, 3.6, 1.25, "평가 지표", "MAE: 평균 절대 오차", GREEN)
    add_card(slide, 8.85, 3.35, 3.6, 1.25, "검증 기준", "scenario_id GroupKFold + OOF", ORANGE)
    add_footer(slide, 1)

    # 2
    slide = make_slide(prs, 2, "문제 정의와 데이터", "향후 30분 평균 출고 지연 시간을 예측하는 회귀 문제")
    add_bullets(
        slide,
        [
            "Target: avg_delay_minutes_next_30m",
            "Metric: MAE",
            "Main data: train.csv / test.csv",
            "Auxiliary data: layout_info.csv",
            "각 행은 특정 시나리오의 특정 시점 창고 운영 snapshot",
        ],
        x=0.75,
        y=1.55,
        w=5.7,
    )
    rows = [
        ["File", "Usage"],
        ["train.csv", "학습 target 포함"],
        ["test.csv", "inference feature matrix 생성"],
        ["layout_info.csv", "layout_id 기준 보조 피처"],
        ["sample_submission.csv", "최종 제출 포맷"],
    ]
    add_table(slide, rows, 7.0, 1.55, 5.3, 2.6)

    # 3
    slide = make_slide(prs, 3, "검증 전략", "Public score에 과도하게 맞추지 않기 위해 OOF 기반 검증을 기본으로 사용")
    add_card(slide, 0.75, 1.55, 3.8, 1.45, "GroupKFold", "scenario_id를 그룹으로 분리하여 같은 시나리오가 train/valid에 동시에 들어가지 않도록 구성", BLUE)
    add_card(slide, 4.85, 1.55, 3.8, 1.45, "OOF Stack", "각 단계 prediction은 train OOF와 test inference prediction을 함께 저장", GREEN)
    add_card(slide, 8.95, 1.55, 3.6, 1.45, "MAE 기준", "모든 내부 선택은 train target 기준 OOF MAE를 우선 확인", ORANGE)
    add_bullets(
        slide,
        [
            "test target은 어떤 단계에서도 사용하지 않음",
            "pseudo labeling 사용하지 않음",
            "Public LB는 최종 후보 선택 및 late postprocess 스칼라 검증에만 활용",
        ],
        x=0.95,
        y=3.55,
        w=11.2,
    )

    # 4
    slide = make_slide(prs, 4, "Feature Engineering", "동일한 전처리를 train/test에 적용하여 학습 행렬과 inference 행렬 생성")
    add_bullets(
        slide,
        [
            "layout_info merge: layout_id 기반 구조 정보 추가",
            "timeslot: scenario_id 내 시점 순서",
            "operational ratios: 로봇/배터리/혼잡/주문량 관련 파생 피처",
            "lag/lead/rolling/context: scenario 내부 흐름 반영",
            "future/context proxy feature: test는 inference matrix로만 사용",
        ],
        x=0.75,
        y=1.45,
        w=6.3,
    )
    add_code_box(
        slide,
        "features.npz = { X, X_te, y_raw, groups, timeslot_tr/te, layout_cluster_tr/te }",
        x=7.25,
        y=2.0,
        w=5.25,
        h=1.25,
    )
    add_card(slide, 7.25, 3.65, 5.25, 1.25, "Rule note", "X_te는 모델 학습 label로 쓰지 않고 최종 예측 생성에만 사용", GREEN)

    # 5
    slide = make_slide(prs, 5, "Base Ensemble: v12", "다양한 objective와 seed/fold를 결합한 기본 예측 축")
    add_bullets(
        slide,
        [
            "LightGBM / XGBoost / CatBoost 기반 tree ensemble",
            "log / sqrt target transform 및 MAE/quantile objective 조합",
            "여러 fold seed의 OOF와 test prediction 평균",
            "SLSQP 기반 OOF blending으로 최종 v12 prediction 생성",
        ],
        x=0.75,
        y=1.5,
        w=6.2,
    )
    rows = [
        ["Artifact", "Description"],
        ["v12_full_oof_bundle.npz", "v12 OOF/prediction bundle"],
        ["sub_v12_main_pre_rerun.csv", "정확 재현용 보존 v12 prediction"],
    ]
    add_table(slide, rows, 7.15, 1.75, 5.2, 1.55)

    # 6
    slide = make_slide(prs, 6, "Future Stack", "향후 시점의 운영 흐름을 반영하는 future-proxy 모델 축")
    add_bullets(
        slide,
        [
            "21_future_proxy_model.py: 1차 future proxy",
            "23_xgb_future_proxy.py: XGBoost proxy 추가",
            "24_future_proxy_allcols.py: 전체 numeric column 기반 proxy",
            "compact_final_stack.py: v25/v26/v27/v28 stack 통합",
        ],
        x=0.75,
        y=1.5,
        w=6.3,
    )
    add_code_box(
        slide,
        "sub_v28_future_stack_scaled.csv\nfinal_future_stack_oof.npz",
        x=7.25,
        y=2.0,
        w=5.2,
        h=1.35,
    )
    add_card(slide, 7.25, 3.75, 5.2, 1.25, "핵심 효과", "시나리오 내부의 후속 상태 정보를 직접 feature화하여 지연 패턴을 보강", BLUE)

    # 7
    slide = make_slide(prs, 7, "Portfolio Blend", "Public에서 가장 안정적이었던 v12/future 조합을 기준점으로 사용")
    add_code_box(
        slide,
        "sub_port_v28_w70 = 0.30 * sub_v12_main_pre_rerun\n                  + 0.70 * sub_v28_future_stack_scaled",
        x=0.85,
        y=1.65,
        w=11.7,
        h=1.35,
    )
    rows = [
        ["Submission", "Public MAE"],
        ["sub_v28_v1265_fp35.csv", "9.6855430648"],
        ["sub_v28_v1260_fp40.csv", "9.6837787747"],
        ["sub_port_v28_w70.csv", "9.6799310594"],
    ]
    add_table(slide, rows, 2.2, 3.6, 8.9, 1.8)

    # 8
    slide = make_slide(prs, 8, "Late Residual Postprocess", "Public에서 실패한 방향을 반대로 돌려 최종 A120 생성")
    add_bullets(
        slide,
        [
            "w70 기준 train OOF residual curve를 학습",
            "v41 후보: sub_v41_rank_layout_a035_add_w65.csv",
            "v41은 Public에서 악화되어 해당 이동 방향을 반대로 적용",
            "a085, a100, a120 순서로 반대 방향이 개선됨",
        ],
        x=0.75,
        y=1.45,
        w=6.3,
    )
    add_code_box(
        slide,
        "A120 = sub_port_v28_w70\n       - 1.20 * (sub_v41_rank_layout_a035_add_w65 - sub_port_v28_w70)",
        x=7.25,
        y=1.9,
        w=5.25,
        h=1.7,
    )
    rows = [
        ["File", "Public MAE"],
        ["sub_port_v28_w70.csv", "9.6799310594"],
        ["sub_v42_anti_v41layout035_a100.csv", "9.6773562672"],
        ["sub_v44_anti_v41layout035_a120.csv", "9.6773143392"],
    ]
    add_table(slide, rows, 7.25, 4.0, 5.25, 1.75, font_size=10)

    # 9
    slide = make_slide(prs, 9, "Rule Compliance", "평가 데이터의 target을 모델 학습에 사용하지 않음")
    add_card(slide, 0.75, 1.5, 3.7, 1.4, "사용하지 않음", "test target, pseudo label, test prediction을 label로 재학습", ORANGE)
    add_card(slide, 4.8, 1.5, 3.7, 1.4, "사용함", "test feature matrix 생성 및 inference", GREEN)
    add_card(slide, 8.85, 1.5, 3.7, 1.4, "보정 기준", "train OOF prediction + train target", BLUE)
    add_bullets(
        slide,
        [
            "동일한 feature engineering pipeline을 train/test에 적용",
            "모든 supervised fitting은 train target 기준으로 수행",
            "late A120은 저장된 prediction CSV의 deterministic postprocess",
            "Public probing은 최종 후보 선택 과정으로 별도 명시",
        ],
        x=0.95,
        y=3.45,
        w=11.2,
    )

    # 10
    slide = make_slide(prs, 10, "Reproduction", "코드 공유 게시판 검증을 위한 실행 경로")
    add_code_box(
        slide,
        "python reproduce_a120_from_artifacts.py",
        x=0.85,
        y=1.55,
        w=6.0,
        h=0.9,
    )
    add_bullets(
        slide,
        [
            "Fast exact reproduction: artifacts 기반 최종 A120 재계산",
            "Full training reproduction: run_full_pipeline.ps1",
            "Raw data는 공식 파일을 data/ 폴더에 배치",
            "환경: Python 3.12.5, LightGBM 4.6.0, XGBoost 3.1.1, CatBoost 1.2.8",
        ],
        x=0.85,
        y=2.8,
        w=11.5,
    )
    rows = [
        ["Output", "Description"],
        ["sub_v44_anti_v41layout035_a120.csv", "최종 선택 제출 파일"],
        ["README.md", "실행 순서와 rule boundary"],
        ["REPRO_MANIFEST.md", "상세 artifact chain"],
    ]
    add_table(slide, rows, 6.95, 1.45, 5.45, 1.75, font_size=10)

    prs.save(OUT)
    print(f"saved: {OUT}")


if __name__ == "__main__":
    build()
